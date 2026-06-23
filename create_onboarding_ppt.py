from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

slides = [
    ("Onboarding Wizard", "Purpose • Your Name • Date", "Open with a 15–30s summary: why onboarding matters (time-to-value, compliance)."),
    ("Agenda", "Overview, Problem, Solution, Flow, Tech, Metrics, Demo, Next steps, Q&A", "Briefly run through agenda so panel knows the structure."),
    ("Problem Statement", "Manual onboarding is slow, error-prone, inconsistent; abandoned signups and compliance risk.", "Give 2 quick examples (delays, missing KYC)."),
    ("Objectives", "Reduce completion time • Improve accuracy • Ensure compliance • Increase activation rate", "Link objectives to business KPIs (activation, support load)."),
    ("What the Onboarding Wizard Is", "Guided multi-step flow that collects user data, validates KYC, uploads documents, and provisions accounts.", "Emphasize it's modular, role-aware, and configurable."),
    ("User Journey (Flow)", "1 Welcome → 2 Basic details → 3 KYC docs → 4 Review & submit → 5 Confirmation", "Walk panel through typical user screen transitions."),
    ("Key Features", "Inline validation • Resumable sessions • Document uploads • Conditional steps • Progress saving • Role-based paths", "Call out resumable sessions and conditional logic as high-impact."),
    ("UX Highlights", "Clear progress bar • Contextual help • Mobile-first • Minimal fields per step • Accessibility", "Mention accessibility and localized strings already present in the repo."),
    ("Technical Architecture", "Frontend (Razor/JS) • Backend APIs • Dapper data layer • File storage for uploads • SQL scripts", "Keep high-level; be ready to point to implementation files if asked."),
    ("Security & Compliance", "Encryption in transit/at rest • Retention policies • RBAC • Audit trails", "Reassure panel on KYC handling and PII exposure in logs."),
    ("Metrics to Track", "Time-to-complete • Abandonment rate • Activations • Support tickets • Document rejection rate", "Suggest baseline and target improvements (e.g., reduce time-to-complete by 40%)."),
    ("Demo Plan", "3–4 minute live demo: create account, upload KYC, resumable session, admin review. Backup: recording.", "State fallback if live demo fails (recording + screenshots)."),
    ("Roadmap & Next Steps", "Improve analytics • Add conditional paths • Integrate ID verification • Automate remediation", "Offer phased timeline (MVP → enhancements → integrations)."),
    ("Risks & Mitigations", "Risk: document fraud → Mitigation: third-party verification; Risk: high abandonment → Mitigation: reduce fields/add help.", "Be concise and propose mitigations."),
    ("Call to Action / Ask", "Decide pilot scope • Select user group • Allocate QA/dev time • Approve vendor budget", "End with specific asks and timeline for pilot."),
    ("Q&A", "Invite questions; appendix ready for deep dive (data model, SQL, APIs).", "Offer to show code/files on request.")
]

for title, content, notes in slides:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    title_placeholder = slide.shapes.title
    body = slide.shapes.placeholders[1]
    title_placeholder.text = title
    tf = body.text_frame
    tf.text = content

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes

# Save
out = "Onboarding_Wizard.pptx"
prs.save(out)
print(f"Saved {out}")
